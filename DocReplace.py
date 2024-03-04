import tkinter as tk
from tkinter import filedialog
from tkinter import messagebox
from docx import Document
from openpyxl import load_workbook
import os

from pptx import Presentation

def batch_replace_text(folder_path):
    old_text = old_text_entry.get()
    new_text = new_text_entry.get()
    filetypes = (".docx", ".txt", ".xlsx", ".pptx")  # 添加支持的文件类型，包括PPT
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            if file.endswith(filetypes):
                filepath = os.path.join(root, file)
                if file.endswith(".docx") or file.endswith(".txt"):
                    replace_text(filepath, old_text, new_text)
                elif file.endswith(".xlsx"):
                    replace_text_in_excel(filepath, old_text, new_text)
                elif file.endswith(".pptx"):
                    replace_text_in_ppt(filepath, old_text, new_text)
    messagebox.showinfo("Success", "Batch text replacement completed!")

def browse_files():
    filetypes = (("Word files", "*.docx"), ("Text files", "*.txt"), ("Excel files", "*.xlsx"), ("PPT files", "*.pptx"), ("All files", "*.*"))
    filepaths = filedialog.askopenfilenames(filetypes=filetypes)
    if filepaths:
        old_text = old_text_entry.get()
        new_text = new_text_entry.get()
        for filepath in filepaths:
            if filepath.endswith(".docx") or filepath.endswith(".txt"):
                replace_text(filepath, old_text, new_text)
            elif filepath.endswith(".xlsx"):
                replace_text_in_excel(filepath, old_text, new_text)
            elif filepath.endswith(".pptx"):
                replace_text_in_ppt(filepath, old_text, new_text)

def replace_text(filepath, old_text, new_text):
    try:
        if filepath.endswith(".docx"):
            document = Document(filepath)
            for paragraph in document.paragraphs:
                paragraph.text = paragraph.text.replace(old_text, new_text)
            document.save(filepath)
        elif filepath.endswith(".txt"):
            with open(filepath, "r") as file:
                content = file.read()
            content = content.replace(old_text, new_text)
            with open(filepath, "w") as file:
                file.write(content)
        messagebox.showinfo("Success", "Text replacement completed!")
    except Exception as e:
        messagebox.showerror("Error", str(e))

def replace_text_in_excel(filepath, old_text, new_text):
    try:
        workbook = load_workbook(filepath)
        for sheetname in workbook.sheetnames:
            sheet = workbook[sheetname]
            for row in sheet.iter_rows():
                for cell in row:
                    if cell.value:
                        if isinstance(cell.value, str):
                            cell.value = cell.value.replace(old_text, new_text)
                        else:
                            cell.value = str(cell.value).replace(old_text, new_text)
        workbook.save(filepath)
        messagebox.showinfo("Success", "Text replacement completed!")
    except Exception as e:
        messagebox.showerror("Error", str(e))


def replace_text_in_ppt(filepath, old_text, new_text):
    try:
        presentation = Presentation(filepath)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            run.text = run.text.replace(old_text, new_text)
        presentation.save(filepath)
        messagebox.showinfo("Success", "Text replacement completed!")
    except Exception as e:
        messagebox.showerror("Error", str(e))

def perform_text_replacement():
    folder_path = filedialog.askdirectory()
    if folder_path:
        batch_replace_text(folder_path)

# 创建主窗口
window = tk.Tk()
window.title("文本处理")
window.geometry("500x300")

# 创建浏览按钮
browse_button = tk.Button(window, text="浏览文件", command=browse_files)
browse_button.pack(pady=10)

# 创建旧文本输入框
old_text_label = tk.Label(window, text="旧文本:", font=("Arial", 12))
old_text_label.pack()
old_text_entry = tk.Entry(window, font=("Arial", 12))
old_text_entry.pack()

# 创建目的文本输入框
new_text_label = tk.Label(window, text="目的文本:", font=("Arial", 12))
new_text_label.pack()
new_text_entry = tk.Entry(window, font=("Arial", 12))
new_text_entry.pack()

# 运行主循环
window.mainloop()
